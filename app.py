import time
import io
import requests
import streamlit as st
import pandas as pd
import numpy as np
from pptx import Presentation
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import LabelEncoder, StandardScaler
from sklearn.neighbors import KNeighborsClassifier
from sklearn.tree import DecisionTreeClassifier, DecisionTreeRegressor
from sklearn.ensemble import RandomForestClassifier, GradientBoostingClassifier
from sklearn.linear_model import LinearRegression, Ridge, Lasso
from sklearn.metrics import (
    accuracy_score, precision_score, recall_score, f1_score,
    confusion_matrix, roc_curve
)
from sklearn.cluster import KMeans
import plotly.express as px
import plotly.graph_objects as go
from mlxtend.frequent_patterns import apriori, association_rules

# ----------------- PAGE CONFIG -----------------
st.set_page_config(page_title="EduNest - EdTech Analytics Dashboard", layout="wide")
start_time = time.time()

# --------------- CUSTOM CSS & ANIMATIONS ---------------
st.markdown(
    """
    <style>
    body { background-color: #1e222e; color: #eee; }
    .stButton>button { border-radius:12px; box-shadow:2px 2px 5px rgba(0,0,0,0.5); }
    </style>
    """,
    unsafe_allow_html=True
)

# ---------- LOGO, TITLE, CAPTION ----------
logo_path = "logo.png"
col_logo, col_title = st.columns([1, 7])
with col_logo:
    st.image(logo_path, width=100)
with col_title:
    st.title("EduNest Survey Insights Dashboard")
    st.caption("EduNest • Unlocking the Power of Learner Analytics")
st.markdown("---")

# ----------- DATA LOADING & QUALITY ------------
@st.cache_data
def load_data(path):
    try:
        return pd.read_csv(path)
    except:
        return None

st.sidebar.header("🔗 Data Source")
uploaded = st.sidebar.file_uploader("Upload survey CSV", type="csv", help="Upload a new CSV to refresh data")
if uploaded:
    df = pd.read_csv(uploaded)
    st.sidebar.success("Using uploaded data")
else:
    df = load_data("data/EduTech_Survey_Synthetic.csv")
    if df is None:
        st.sidebar.error("No data found! Please upload a CSV.")
        st.stop()
    else:
        st.sidebar.info("Using default data from repo")

# Data quality & automated alerts
n_missing = df.isnull().sum().sum()
numeric = df.select_dtypes(include=[np.number])
n_outliers = int(((numeric - numeric.mean()).abs() / numeric.std() > 4).sum().sum())
sub_rate = 100 * df['Paid Subscription'].value_counts(normalize=True).get('Yes', 0)
threshold_sub = 30
if n_missing or n_outliers:
    st.warning(f"⚠️ Data Quality: {n_missing} missing values, {n_outliers} outliers detected.")
else:
    st.success("✅ Data Quality: Clean")

if sub_rate < threshold_sub:
    st.error(f"🚨 Alert: Subscription rate is only {sub_rate:.1f}% (< {threshold_sub}%)!")
else:
    st.info(f"Subscription rate: {sub_rate:.1f}%")

# ------------ SIDEBAR FILTERS -------------
with st.sidebar.expander("🧲 Filter Data", expanded=False):
    st.write("Drill down by demographics")
    sel_age = st.multiselect("Age", sorted(df['Age'].unique()), default=sorted(df['Age'].unique()))
    sel_gen = st.multiselect("Gender", sorted(df['Gender'].unique()), default=sorted(df['Gender'].unique()))
    sel_inc = st.multiselect("Monthly Income", sorted(df['Monthly Income'].unique()), default=sorted(df['Monthly Income'].unique()))
    df = df[df['Age'].isin(sel_age) & df['Gender'].isin(sel_gen) & df['Monthly Income'].isin(sel_inc)]

# ---------------- TABS ----------------
tabs = st.tabs([
    "Executive Summary",
    "Data Visualisation",
    "Time-Series Analysis",
    "Classification",
    "Clustering",
    "Association Rules",
    "Regression",
    "Download PPT",
    "Feedback"
])

# ===== 1. EXECUTIVE SUMMARY =====
with tabs[0]:
    st.subheader("📋 Executive Summary")
    st.info(f"""
- **Subscription Rate:** {sub_rate:.1f}%  
- **Avg Satisfaction:** {df['Satisfaction Level'].mean():.2f}/5  
- **Avg Comfort:** {df['App Comfort Level'].mean():.2f}/5  
- **Top Device:** {df['Device'].mode()[0]}  
- **Avg Free Trial Willingness:** {df['Free Trial Willingness'].mean():.2f}/5  
""")

# ===== 2. DATA VISUALISATION =====
with tabs[1]:
    st.header("📊 Data Visualisation")
    st.info("Explore interactive charts. Use the theme toggle and filters.")
    theme = st.radio("Theme", ["Light", "Dark"], horizontal=True)
    tpl = "plotly_white" if theme == "Light" else "plotly_dark"

    with st.expander("📖 How to Use"):
        st.write("""
- Sidebar filters apply globally.  
- Hover for details.  
- Download data or PPT summary.  
""")

    # KPI Cards
    k1, k2, k3, k4 = st.columns(4)
    k1.metric("Subscription Rate", f"{sub_rate:.1f}%")
    k2.metric("Avg Satisfaction", f"{df['Satisfaction Level'].mean():.2f}/5")
    k3.metric("Avg Comfort", f"{df['App Comfort Level'].mean():.2f}/5")
    k4.metric("Avg Free Trial", f"{df['Free Trial Willingness'].mean():.2f}/5")

    # Plots
    c1, c2 = st.columns(2)
    figs = [
        px.histogram(df, x="Age", color="Age", title="Age Distribution", template=tpl),
        px.pie(df, names="Gender", title="Gender Distribution", template=tpl),
        px.histogram(df, x="Monthly Income", color="Paid Subscription", barmode="group", title="Income vs Subscription", template=tpl),
        px.bar(df['Device'].value_counts().reset_index(name='Count').rename(columns={'index':'Device'}), x="Device", y="Count", title="Preferred Devices", template=tpl)
    ]
    for i, fig in enumerate(figs):
        (c1 if i % 2 == 0 else c2).plotly_chart(fig, use_container_width=True)

    st.plotly_chart(px.pie(df, names="Internet Quality", title="Internet Quality", template=tpl), use_container_width=True)
    st.plotly_chart(px.bar(df, x="Age", color="Paid Subscription", title="Subscription by Age", template=tpl), use_container_width=True)
    st.plotly_chart(px.bar(df.groupby("Occupation")["Satisfaction Level"].mean().reset_index(), x="Occupation", y="Satisfaction Level", title="Satisfaction by Occupation", template=tpl), use_container_width=True)
    st.plotly_chart(go.Figure(go.Funnel(y=["Exposed", "Signed Up", "Subscribed"], x=[len(df), df['Signed Up'].value_counts().get('Yes', 0), df['Paid Subscription'].value_counts().get('Yes', 0)])), use_container_width=True)
    st.plotly_chart(px.choropleth(df['Country'].value_counts().reset_index(name='Count').rename(columns={'index':'Country'}), locations='Country', locationmode='country names', color='Count', title="Users by Country", template=tpl), use_container_width=True)

    # Correlation heatmap
    num_df = df.select_dtypes(include=[np.number]).copy()
    corr = num_df.corr()
    st.plotly_chart(px.imshow(corr, text_auto=True, color_continuous_scale="RdBu_r", title="Correlation Heatmap", template=tpl), use_container_width=True)

    st.markdown("---")
    st.download_button("Download Filtered Data", df.to_csv(index=False), "filtered_data.csv", "text/csv")

# ===== 3. TIME-SERIES ANALYSIS =====
with tabs[2]:
    st.header("⏱️ Time-Series Analysis")
    if 'Timestamp' in df.columns:
        df['Timestamp'] = pd.to_datetime(df['Timestamp'])
        ts_df = df.set_index('Timestamp').resample('D').agg({
            'Signed Up': lambda x: (x == "Yes").sum(),
            'Paid Subscription': lambda x: (x == "Yes").sum()
        }).rename(columns={'Signed Up': 'Daily Signups', 'Paid Subscription': 'Daily Subscriptions'})
        fig_ts = px.line(ts_df, x=ts_df.index, y=['Daily Signups', 'Daily Subscriptions'], title="Daily Signups & Subscriptions", template=tpl)
        st.plotly_chart(fig_ts, use_container_width=True)
    else:
        st.warning("No 'Timestamp' column found for time-series analysis.")

# ===== 4. CLASSIFICATION =====
with tabs[3]:
    st.header("🧩 Classification")
    st.info("Predict Paid Subscription. Upload new data to predict.")
    clf = df.copy().fillna("Unknown")
    drops = ["Country", "City", "Non-Subscription Reasons", "Features Used Most", "Other Apps Used", "Motivation", "Learning Goals", "Preferred Subjects", "Preferred App Features", "Selection Factors", "Learning Challenges"]
    clf.drop(columns=[c for c in drops if c in clf], inplace=True)
    encoders = {}
    for col in clf.select_dtypes(include='object'):
        if col != "Paid Subscription":
            le = LabelEncoder()
            clf[col] = le.fit_transform(clf[col].astype(str))
            encoders[col] = le
    clf["Paid Subscription"] = (clf["Paid Subscription"] == "Yes").astype(int)
    X = clf.drop("Paid Subscription", axis=1)
    y = clf["Paid Subscription"]
    X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.25, random_state=42)
    scaler = StandardScaler()
    X_train_s, X_test_s = scaler.fit_transform(X_train), scaler.transform(X_test)
    models = {
        "KNN": KNeighborsClassifier(),
        "Decision Tree": DecisionTreeClassifier(max_depth=6, random_state=0),
        "Random Forest": RandomForestClassifier(n_estimators=60, random_state=0),
        "GBRT": GradientBoostingClassifier(n_estimators=60, random_state=0)
    }
    results = {}
    probs = {}
    for name, model in models.items():
        model.fit(X_train_s, y_train)
        preds = model.predict(X_test_s)
        prob = model.predict_proba(X_test_s)[:, 1]
        results[name] = {
            "Train Acc": accuracy_score(y_train, model.predict(X_train_s)),
            "Test Acc": accuracy_score(y_test, preds),
            "Precision": precision_score(y_test, preds),
            "Recall": recall_score(y_test, preds),
            "F1": f1_score(y_test, preds)
        }
        probs[name] = prob
    st.dataframe(pd.DataFrame(results).T.round(3))
    sel = st.selectbox("Confusion Matrix", list(models.keys()))
    cm = confusion_matrix(y_test, models[sel].predict(X_test_s))
    st.plotly_chart(px.imshow(cm, text_auto=True, labels=dict(x="Pred", y="True"), title=f"{sel} Confusion Matrix", template=tpl), use_container_width=True)
    roc_fig = go.Figure()
    for name, p in probs.items():
        fpr, tpr, _ = roc_curve(y_test, p)
        roc_fig.add_trace(go.Scatter(x=fpr, y=tpr, mode='lines', name=name))
    roc_fig.add_trace(go.Scatter(x=[0,1], y=[0,1], mode='lines', name='Random', line=dict(dash='dash')))
    st.plotly_chart(roc_fig, use_container_width=True)
    up = st.file_uploader("Upload CSV to Predict", type="csv", key="pred_up")
    if up:
        pdf = pd.read_csv(up)
        for col, le in encoders.items():
            if col in pdf:
                pdf[col] = le.transform(pdf[col].astype(str))
        p_s = scaler.transform(pdf[X.columns])
        pdf["Predicted Subscription"] = np.where(models[sel].predict(p_s) == 1, "Yes", "No")
        st.dataframe(pdf.head())
        st.download_button("Download Predictions", pdf.to_csv(index=False), "predictions.csv", "text/csv")

# ===== 5. CLUSTERING =====
with tabs[4]:
    st.header("👥 Clustering")
    st.info("Segment customers. Adjust number of clusters.")
    cdf = df.copy()
    cols = ['Age', 'Gender', 'Monthly Income', 'Occupation', 'Courses Taken Last Year', 'Device']
    for c in cols:
        cdf[c] = LabelEncoder().fit_transform(cdf[c].astype(str))
    k = st.slider("Number of clusters", 2, 10, 4)
    km = KMeans(n_clusters=k, random_state=0, n_init=10).fit(cdf[cols])
    cdf['Cluster'] = km.labels_
    inertias = [KMeans(n_clusters=i, random_state=0, n_init=10).fit(cdf[cols]).inertia_ for i in range(2, 11)]
    st.plotly_chart(px.line(x=list(range(2, 11)), y=inertias, title="Elbow Chart", markers=True, template=tpl), use_container_width=True)
    centers = pd.DataFrame(km.cluster_centers_, columns=cols)
    st.dataframe(centers.round(2))
    st.download_button("Download Clustered Data", cdf.to_csv(index=False), "clustered.csv", "text/csv")

# ===== 6. ASSOCIATION RULES =====
with tabs[5]:
    st.header("🔗 Association Rules")
    st.info("Mine top-10 rules from multi-select questions.")
    multi_cols = ['Motivation', 'Learning Goals', 'Preferred Subjects', 'Preferred App Features',
                  'Selection Factors', 'Learning Challenges', 'Non-Subscription Reasons',
                  'Features Used Most', 'Other Apps Used']
    sel_cols = st.multiselect("Select columns", multi_cols, default=multi_cols[:2])
    sup = st.slider("Min Support", 0.01, 0.2, 0.05)
    conf = st.slider("Min Confidence", 0.2, 1.0, 0.6)
    if len(sel_cols) >= 2:
        transactions = []
        for _, row in df[sel_cols].iterrows():
            items = set()
            for col in sel_cols:
                for v in str(row[col]).split(','):
                    vv = v.strip()
                    if vv and vv.lower() != 'none':
                        items.add(f"{col}:{vv}")
            transactions.append(list(items))
        from mlxtend.preprocessing import TransactionEncoder
        te = TransactionEncoder().fit(transactions)
        df_te = pd.DataFrame(te.transform(transactions), columns=te.columns_)
        freq = apriori(df_te, min_support=sup, use_colnames=True)
        rules = association_rules(freq, metric="confidence", min_threshold=conf).sort_values('confidence', ascending=False).head(10)
        rules['antecedents'] = rules['antecedents'].apply(lambda x: ", ".join(x))
        rules['consequents'] = rules['consequents'].apply(lambda x: ", ".join(x))
        st.dataframe(rules[['antecedents', 'consequents', 'support', 'confidence', 'lift']])
    else:
        st.warning("Select at least 2 columns to run association rules.")

# ===== 7. REGRESSION =====
with tabs[6]:
    st.header("📈 Regression")
    st.info("Run and compare regression models.")
    targets = {"Spend": "Willingness to Spend", "Satisfaction": "Satisfaction Level"}
    choice = st.selectbox("Select target", list(targets.keys()))
    ycol = targets[choice]
    if ycol == "Willingness to Spend":
        y = df[ycol].map({'<10': 1, '10-20': 2, '20-50': 3, '50-100': 4, '>100': 5})
    else:
        y = df[ycol]
    X = df[['Age', 'Gender', 'Education', 'Monthly Income', 'Occupation', 'Device',
            'Internet Quality', 'Courses Taken Last Year', 'App Comfort Level']].copy()
    for c in X.select_dtypes(include='object'):
        X[c] = LabelEncoder().fit_transform(X[c].astype(str))
    X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.25, random_state=42)
    sc = StandardScaler()
    X_train_s, X_test_s = sc.fit_transform(X_train), sc.transform(X_test)
    regmods = {
        "Linear": LinearRegression(),
        "Ridge": Ridge(),
        "Lasso": Lasso(alpha=0.01),
        "Tree": DecisionTreeRegressor(max_depth=6, random_state=0)
    }
    out, preds = {}, {}
    for name, model in regmods.items():
        model.fit(X_train_s, y_train)
        p = model.predict(X_test_s)
        out[name] = {"Train R2": model.score(X_train_s, y_train), "Test R2": model.score(X_test_s, y_test)}
        preds[name] = p
    st.dataframe(pd.DataFrame(out).T.round(3))
    for name, p in preds.items():
        st.plotly_chart(px.scatter(x=y_test, y=p, labels={'x': 'Actual', 'y': 'Predicted'},
                                   title=f"{name} Actual vs Predicted", template=tpl), use_container_width=True)
    imp = regmods["Tree"].feature_importances_
    st.plotly_chart(px.bar(x=X.columns, y=imp, title="Tree Feature Importances", template=tpl), use_container_width=True)

# ===== 8. DOWNLOAD PPT SUMMARY =====
with tabs[7]:
    st.header("📥 Download PPT Summary")
    bullets = [
        f"Subscription Rate: {sub_rate:.1f}%",
        f"Avg Satisfaction: {df['Satisfaction Level'].mean():.2f}/5",
        f"Avg Comfort Level: {df['App Comfort Level'].mean():.2f}/5",
        f"Top Device Used: {df['Device'].mode()[0]}"
    ]
    if st.button("Generate & Download PPT"):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "EduNest Executive Summary"
        tf = slide.shapes.placeholders[1].text_frame
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.level = 1
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        st.download_button("Download PPT", ppt_io, "EduNest_Summary.pptx",
                           "application/vnd.openxmlformats-officedocument.presentationml.presentation")

# ===== 9. FEEDBACK =====
with tabs[8]:
    st.header("💬 Feedback")
    fb = st.text_area("Enter your suggestions or comments:")
    if st.button("Submit Feedback"):
        st.success("Thank you for your feedback!")
        webhook = st.secrets.get("slack_webhook_url")
        if webhook and fb:
            try:
                requests.post(webhook, json={"text": f"New feedback: {fb}"})
            except:
                st.warning("Failed to send feedback to Slack.")

# ---- FINAL TAGLINE & PERFORMANCE ----
st.markdown("---")
st.markdown(
    "<div style='text-align:center; color:#777; font-size:14px;'>"
    "MBA Data Analytics for Insights & Decision Making Project • Team Subhayu | Unlocking end-to-end learner insights"
    "</div>",
    unsafe_allow_html=True
)
st.caption(f"Dashboard generated in {time.time() - start_time:.2f}s")
